"""Generate synthetic fixtures covering every format family BrainForgeMD claims.

All fixtures are generated locally from code. No user data, no downloaded corpora.

Usage:  python audit/gen_fixtures.py <output-dir>
"""

from __future__ import annotations

import math
import struct
import subprocess
import sys
import wave
import zipfile
from pathlib import Path

OUT = Path(sys.argv[1] if len(sys.argv) > 1 else "audit/fixtures")
OUT.mkdir(parents=True, exist_ok=True)

MADE: list[tuple[str, int]] = []
FAILED: list[tuple[str, str]] = []
UNI = "Cafe — Quebec — 日本語 — العربية — Omega"
FF = "ffmpeg"


def mk(name, fn):
    p = OUT / name
    try:
        fn(p)
        MADE.append((name, p.stat().st_size))
        return p
    except Exception as exc:
        FAILED.append((name, f"{type(exc).__name__}: {exc}"))
        return None


# --------------------------------------------------------------- PDF
def pdf_multipage(p):
    from reportlab.lib.pagesizes import A4
    from reportlab.pdfgen import canvas

    c = canvas.Canvas(str(p), pagesize=A4)
    for i in range(1, 6):
        c.drawString(72, 760, f"PDF page {i} heading")
        c.drawString(72, 730, f"Body content page {i}: MARKER_PAGE_{i}")
        c.showPage()
    c.save()


def pdf_unicode(p):
    from reportlab.pdfbase import pdfmetrics
    from reportlab.pdfbase.ttfonts import TTFont
    from reportlab.pdfgen import canvas

    c = canvas.Canvas(str(p))
    font = Path("C:/Windows/Fonts/arial.ttf")
    if font.exists():
        pdfmetrics.registerFont(TTFont("Arial", str(font)))
        c.setFont("Arial", 14)
        c.drawString(72, 760, "Accents: eaui - Cafe - Quebec - Omega")
    c.drawString(72, 730, "MARKER_UNICODE_PDF")
    c.save()


def pdf_table(p):
    from reportlab.lib.styles import getSampleStyleSheet
    from reportlab.platypus import Paragraph, SimpleDocTemplate, Table

    styles = getSampleStyleSheet()
    data = [["Region", "Q1", "Q2"], ["North", "100", "200"], ["South", "300", "400"]]
    SimpleDocTemplate(str(p)).build(
        [Paragraph("Quarterly MARKER_TABLE_PDF", styles["Title"]), Table(data)]
    )


def pdf_image(p):
    from PIL import Image, ImageDraw
    from reportlab.lib.utils import ImageReader
    from reportlab.pdfgen import canvas

    img = Image.new("RGB", (300, 120), "white")
    ImageDraw.Draw(img).text((10, 50), "IMAGE INSIDE PDF", fill="black")
    c = canvas.Canvas(str(p))
    c.drawString(72, 760, "MARKER_PDF_WITH_IMAGE")
    c.drawImage(ImageReader(img), 72, 600, width=300, height=120)
    c.save()


def pdf_scanned(p):
    """Image-only PDF: text exists only as pixels, so it needs real OCR."""
    from PIL import Image, ImageDraw, ImageFont
    from reportlab.lib.utils import ImageReader
    from reportlab.pdfgen import canvas

    img = Image.new("RGB", (900, 400), "white")
    draw = ImageDraw.Draw(img)
    try:
        font = ImageFont.truetype("C:/Windows/Fonts/arial.ttf", 40)
    except Exception:
        font = ImageFont.load_default()
    draw.text((20, 180), "SCANNED TEXT ONLY IN PIXELS", fill="black", font=font)
    c = canvas.Canvas(str(p))
    c.drawImage(ImageReader(img), 20, 400, width=560, height=250)
    c.save()


def pdf_nearly_empty(p):
    from reportlab.pdfgen import canvas

    c = canvas.Canvas(str(p))
    c.showPage()
    c.save()


def pdf_large(p):
    from reportlab.pdfgen import canvas

    c = canvas.Canvas(str(p))
    for i in range(300):
        c.drawString(72, 770, f"Large PDF page {i} MARKER_LARGE")
        for j in range(30):
            c.drawString(72, 740 - j * 22, f"line {j} page {i} lorem ipsum dolor sit amet consectetur")
        c.showPage()
    c.save()


def pdf_corrupt(p):
    from reportlab.pdfgen import canvas

    tmp = p.with_name("_ok.pdf")
    c = canvas.Canvas(str(tmp))
    c.drawString(72, 760, "content")
    c.save()
    data = bytearray(tmp.read_bytes())
    mid = len(data) // 2
    del data[mid : mid + 400]
    p.write_bytes(bytes(data))
    tmp.unlink()


def pdf_encrypted(p):
    from reportlab.pdfgen import canvas

    c = canvas.Canvas(str(p), encrypt="s3cret")
    c.drawString(72, 760, "MARKER_ENCRYPTED_PDF")
    c.save()


# --------------------------------------------------------------- DOCX
def docx_rich(p):
    import docx
    from PIL import Image

    d = docx.Document()
    d.add_heading("DOCX MARKER_HEADING", 0)
    d.add_heading("Section one", 1)
    d.add_paragraph("Plain paragraph MARKER_DOCX_BODY.")
    d.add_paragraph("Bullet one", style="List Bullet")
    d.add_paragraph("Bullet two", style="List Bullet")
    d.add_paragraph("Numbered one", style="List Number")
    table = d.add_table(rows=2, cols=3)
    for r, row in enumerate([["A", "B", "MARKER_DOCX_TABLE"], ["1", "2", "3"]]):
        for c, value in enumerate(row):
            table.cell(r, c).text = value
    d.add_paragraph(UNI)
    img = OUT / "_tmp_docx.png"
    Image.new("RGB", (80, 40), "blue").save(img)
    d.add_picture(str(img))
    img.unlink()
    section = d.sections[0]
    section.header.paragraphs[0].text = "HEADER MARKER_DOCX_HEADER"
    section.footer.paragraphs[0].text = "FOOTER MARKER_DOCX_FOOTER"
    d.save(str(p))


def docx_corrupt(p):
    import docx

    tmp = p.with_name("_ok.docx")
    d = docx.Document()
    d.add_paragraph("content")
    d.save(str(tmp))
    data = bytearray(tmp.read_bytes())
    data[200:280] = b"\x00" * 80
    p.write_bytes(bytes(data))
    tmp.unlink()


# --------------------------------------------------------------- XLSX
def xlsx_rich(p):
    import openpyxl

    wb = openpyxl.Workbook()
    s1 = wb.active
    s1.title = "Sheet One"
    s1["A1"], s1["B1"] = "Name", "Value"
    s1["A2"], s1["B2"] = "alpha", 10
    s1["A3"], s1["B3"] = "beta", 20
    s1["B4"] = "=SUM(B2:B3)"
    s1["A5"] = UNI
    s1["D9"] = "after empty cells MARKER_XLSX"
    wb.create_sheet("Second Sheet")["A1"] = "SECOND_SHEET_MARKER"
    big = wb.create_sheet("Big")
    for r in range(1, 501):
        for c in range(1, 11):
            big.cell(r, c, f"r{r}c{c}")
    wb.save(str(p))


# --------------------------------------------------------------- PPTX
def pptx_rich(p):
    from PIL import Image
    from pptx import Presentation
    from pptx.util import Inches

    pr = Presentation()
    s = pr.slides.add_slide(pr.slide_layouts[0])
    s.shapes.title.text = "PPTX MARKER_TITLE"
    s.placeholders[1].text = "Subtitle body text"
    s.notes_slide.notes_text_frame.text = "SPEAKER NOTES MARKER_PPTX_NOTES"
    s2 = pr.slides.add_slide(pr.slide_layouts[1])
    s2.shapes.title.text = "Second slide"
    s2.placeholders[1].text = f"Bullet A\nBullet B\n{UNI}"
    s3 = pr.slides.add_slide(pr.slide_layouts[5])
    table = s3.shapes.add_table(2, 2, Inches(1), Inches(2), Inches(4), Inches(1)).table
    table.cell(0, 0).text = "H1"
    table.cell(0, 1).text = "H2"
    table.cell(1, 0).text = "v1"
    table.cell(1, 1).text = "MARKER_PPTX_TABLE"
    img = OUT / "_tmp_pptx.png"
    Image.new("RGB", (60, 60), "red").save(img)
    s3.shapes.add_picture(str(img), Inches(6), Inches(2), Inches(1), Inches(1))
    img.unlink()
    pr.save(str(p))


# --------------------------------------------------------------- OpenDocument
def odt_doc(p):
    from odf.opendocument import OpenDocumentText
    from odf.text import H, P

    d = OpenDocumentText()
    d.text.addElement(H(outlinelevel=1, text="ODT MARKER_HEADING"))
    d.text.addElement(P(text="ODT body MARKER_ODT_BODY"))
    d.text.addElement(P(text=UNI))
    d.save(str(p))


def ods_sheet(p):
    from odf.opendocument import OpenDocumentSpreadsheet
    from odf.table import Table, TableCell, TableRow
    from odf.text import P

    d = OpenDocumentSpreadsheet()
    table = Table(name="Sheet1")
    for row in [["Name", "Value"], ["alpha", "1"], ["MARKER_ODS", "42"]]:
        tr = TableRow()
        for value in row:
            cell = TableCell(valuetype="string")
            cell.addElement(P(text=value))
            tr.addElement(cell)
        table.addElement(tr)
    d.spreadsheet.addElement(table)
    d.save(str(p))


def odp_deck(p):
    from odf.draw import Frame, Page, TextBox
    from odf.opendocument import OpenDocumentPresentation
    from odf.style import MasterPage, PageLayout
    from odf.text import P

    d = OpenDocumentPresentation()
    layout = PageLayout(name="PL")
    d.automaticstyles.addElement(layout)
    master = MasterPage(name="M", pagelayoutname=layout)
    d.masterstyles.addElement(master)
    page = Page(masterpagename="M")
    frame = Frame(width="10cm", height="2cm", x="1cm", y="1cm")
    box = TextBox()
    box.addElement(P(text="ODP MARKER_SLIDE"))
    frame.addElement(box)
    page.addElement(frame)
    d.presentation.addElement(page)
    d.save(str(p))


# --------------------------------------------------------------- Images
def _img(text, size=(500, 180), bg="white", fg="black", fontsize=28):
    from PIL import Image, ImageDraw, ImageFont

    im = Image.new("RGB", size, bg)
    draw = ImageDraw.Draw(im)
    try:
        font = ImageFont.truetype("C:/Windows/Fonts/arial.ttf", fontsize)
    except Exception:
        font = ImageFont.load_default()
    draw.text((10, size[1] // 3), text, fill=fg, font=font)
    return im


def _blank(p):
    from PIL import Image

    Image.new("RGB", (200, 200), "green").save(p)


# --------------------------------------------------------------- Audio
def wav_tone(p, secs=2.0, freq=440.0, amp=12000, sr=16000):
    with wave.open(str(p), "w") as w:
        w.setnchannels(1)
        w.setsampwidth(2)
        w.setframerate(sr)
        w.writeframes(
            b"".join(
                struct.pack("<h", int(amp * math.sin(2 * math.pi * freq * t / sr)))
                for t in range(int(sr * secs))
            )
        )


def wav_noise(p):
    import random

    random.seed(7)
    sr = 16000
    with wave.open(str(p), "w") as w:
        w.setnchannels(1)
        w.setsampwidth(2)
        w.setframerate(sr)
        w.writeframes(b"".join(struct.pack("<h", random.randint(-8000, 8000)) for _ in range(sr * 2)))


def transcode(src, dst, *extra):
    subprocess.run([FF, "-y", "-loglevel", "error", "-i", str(src), *extra, str(dst)], check=True)


def video(p, with_audio):
    args = [FF, "-y", "-loglevel", "error", "-f", "lavfi", "-i", "testsrc=size=320x240:rate=10:duration=2"]
    if with_audio:
        args += ["-f", "lavfi", "-i", "sine=frequency=440:duration=2", "-shortest"]
    args += ["-pix_fmt", "yuv420p", str(p)]
    subprocess.run(args, check=True)


# --------------------------------------------------------------- EPUB
def epub(p):
    with zipfile.ZipFile(p, "w") as z:
        z.writestr("mimetype", "application/epub+zip")
        z.writestr(
            "META-INF/container.xml",
            '<?xml version="1.0"?><container version="1.0" '
            'xmlns="urn:oasis:names:tc:opendocument:xmlns:container"><rootfiles>'
            '<rootfile full-path="OEBPS/content.opf" '
            'media-type="application/oebps-package+xml"/></rootfiles></container>',
        )
        z.writestr(
            "OEBPS/content.opf",
            '<?xml version="1.0"?><package xmlns="http://www.idpf.org/2007/opf" version="3.0" '
            'unique-identifier="id"><metadata xmlns:dc="http://purl.org/dc/elements/1.1/">'
            '<dc:identifier id="id">x1</dc:identifier><dc:title>Synthetic Book</dc:title>'
            "<dc:language>en</dc:language></metadata><manifest>"
            '<item id="c1" href="ch1.xhtml" media-type="application/xhtml+xml"/></manifest>'
            '<spine><itemref idref="c1"/></spine></package>',
        )
        z.writestr(
            "OEBPS/ch1.xhtml",
            '<?xml version="1.0"?><html xmlns="http://www.w3.org/1999/xhtml"><head>'
            "<title>Chapter</title></head><body><h1>Chapter one</h1>"
            "<p>EPUB MARKER_BODY text.</p></body></html>",
        )


def legacy_xls(p):
    """Genuine legacy BIFF8 .xls binary workbook (not an XLSX rename)."""
    import xlwt

    book = xlwt.Workbook()
    sheet = book.add_sheet("Legacy")
    for r, row in enumerate([["Name", "Value"], ["alpha", "1"], ["MARKER_LEGACY_XLS", "42"]]):
        for c, value in enumerate(row):
            sheet.write(r, c, value)
    book.save(str(p))


def parquet(p):
    import pyarrow as pa
    import pyarrow.parquet as pq

    table = pa.table({"id": list(range(300)), "name": [f"n{i}" for i in range(300)], "note": [UNI] * 300})
    pq.write_table(table, str(p))


def main():
    for name, fn in [
        ("pdf_multipage.pdf", pdf_multipage),
        ("pdf_unicode.pdf", pdf_unicode),
        ("pdf_table.pdf", pdf_table),
        ("pdf_image.pdf", pdf_image),
        ("pdf_scanned.pdf", pdf_scanned),
        ("pdf_nearly_empty.pdf", pdf_nearly_empty),
        ("pdf_large.pdf", pdf_large),
        ("pdf_corrupt.pdf", pdf_corrupt),
        ("pdf_encrypted.pdf", pdf_encrypted),
        ("docx_rich.docx", docx_rich),
        ("docx_corrupt.docx", docx_corrupt),
        ("xlsx_rich.xlsx", xlsx_rich),
        ("pptx_rich.pptx", pptx_rich),
        ("odt_doc.odt", odt_doc),
        ("ods_sheet.ods", ods_sheet),
        ("odp_deck.odp", odp_deck),
        ("book.epub", epub),
        ("data.parquet", parquet),
        ("legacy_book.xls", legacy_xls),
    ]:
        mk(name, fn)

    for name, fn in [
        ("img_text.png", lambda p: _img("OCR MARKER PNG TEXT").save(p)),
        ("img_text.jpg", lambda p: _img("OCR MARKER JPEG TEXT").save(p, quality=92)),
        ("img_text.tiff", lambda p: _img("OCR MARKER TIFF TEXT").save(p)),
        ("img_text.webp", lambda p: _img("OCR MARKER WEBP TEXT").save(p)),
        ("img_text.bmp", lambda p: _img("OCR MARKER BMP TEXT").save(p)),
        ("img_notext.png", _blank),
        ("img_tiny_text.png", lambda p: _img("tiny six point text", size=(300, 60), fontsize=6).save(p)),
        ("img_rotated.png", lambda p: _img("ROTATED OCR MARKER").rotate(90, expand=True).save(p)),
        ("img_lowcontrast.png", lambda p: _img("LOW CONTRAST MARKER", bg="#808080", fg="#8a8a8a").save(p)),
        ("img_corrupt.png", lambda p: p.write_bytes(b"\x89PNG\r\n\x1a\n" + b"\x00" * 100)),
    ]:
        mk(name, fn)

    mk("audio_tone.wav", lambda p: wav_tone(p))
    mk("audio_silence.wav", lambda p: wav_tone(p, amp=0))
    mk("audio_veryshort.wav", lambda p: wav_tone(p, secs=0.05))
    mk("audio_noise.wav", wav_noise)
    tone = OUT / "audio_tone.wav"
    if tone.exists():
        for name in ["audio_tone.mp3", "audio_tone.flac", "audio_tone.ogg", "audio_tone.m4a"]:
            mk(name, lambda p, s=tone: transcode(s, p))

    for name, audio in [
        ("video_noaudio.mp4", False),
        ("video_audio.mp4", True),
        ("video_audio.mkv", True),
        ("video_audio.webm", True),
        ("video_audio.mov", True),
        ("video_audio.avi", True),
    ]:
        mk(name, lambda p, a=audio: video(p, a))

    print(f"GENERATED {len(MADE)} fixtures")
    for name, size in sorted(MADE):
        print(f"  {name:<26} {size:>10} bytes")
    if FAILED:
        print(f"\nCOULD NOT GENERATE {len(FAILED)}:")
        for name, why in sorted(FAILED):
            print(f"  {name:<26} {why[:110]}")


if __name__ == "__main__":
    main()
